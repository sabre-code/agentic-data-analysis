"""
PPTX Generator Service

Creates professional PowerPoint presentations using python-pptx.
Includes title slide, key metrics, chart slides, and insights.
"""
from __future__ import annotations

import logging
from io import BytesIO
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

logger = logging.getLogger(__name__)


class PPTXGenerator:
    """Generates professional PowerPoint presentations from analysis artifacts."""

    def __init__(self) -> None:
        self.slide_width = Inches(10)  # 16:9 aspect ratio
        self.slide_height = Inches(5.625)

    def generate(
        self,
        output_path: Path,
        user_query: str,
        file_schema: dict[str, Any] | None,
        code_result: dict[str, Any] | None,
        code_output: str | None,
        chart_images: list[bytes],
        chart_titles: list[str],
        generated_code: str | None = None,
        executive_summary: str | None = None,
    ) -> None:
        """
        Generate a PPTX presentation.
        
        Args:
            output_path: Path where PPTX will be saved
            user_query: The user's original question
            file_schema: Dataset metadata
            code_result: Structured metrics from code execution
            code_output: stdout from code execution
            chart_images: List of PNG chart images as bytes
            chart_titles: List of titles for each chart
            generated_code: Python code that was executed
            executive_summary: AI-generated executive summary (optional)
        """
        logger.info("Generating PPTX presentation: %s", output_path)
        
        # Create presentation
        prs = Presentation()
        prs.slide_width = self.slide_width
        prs.slide_height = self.slide_height
        
        # Slide 1: Title slide
        self._add_title_slide(prs, file_schema)
        
        # Slide 2: Executive Summary (if provided)
        if executive_summary:
            self._add_executive_summary_slide(prs, executive_summary)
        
        # Slide 3+: Chart slides (one per chart)
        if chart_images:
            for img_bytes, title in zip(chart_images, chart_titles):
                self._add_chart_slide(prs, img_bytes, title)
        
        # Key metrics slide (optional - only if we have structured results)
        if code_result and len(code_result) > 0:
            self._add_metrics_slide(prs, code_result, code_output)
        
        # Save presentation
        prs.save(str(output_path))
        logger.info("PPTX presentation generated successfully: %d bytes", output_path.stat().st_size)

    def _add_title_slide(
        self, 
        prs: Presentation, 
        file_schema: dict[str, Any] | None
    ) -> None:
        """Add title slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2), Inches(9), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = "Data Analysis Report"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(44, 62, 80)  # Dark blue-gray
        title_para.alignment = PP_ALIGN.CENTER
        
        # Dataset info
        if file_schema:
            dataset_name = file_schema.get('original_filename', 'Unknown Dataset')
            row_count = file_schema.get('row_count', 'Unknown')
            
            info_box = slide.shapes.add_textbox(
                Inches(1), Inches(3.5), Inches(8), Inches(0.8)
            )
            info_frame = info_box.text_frame
            info_para = info_frame.paragraphs[0]
            if isinstance(row_count, int):
                info_para.text = f"Dataset: {dataset_name} ({row_count:,} rows)"
            else:
                info_para.text = f"Dataset: {dataset_name} ({row_count} rows)"
            info_para.font.size = Pt(18)
            info_para.font.italic = True
            info_para.font.color.rgb = RGBColor(127, 140, 141)
            info_para.alignment = PP_ALIGN.CENTER

    def _add_executive_summary_slide(
        self,
        prs: Presentation,
        executive_summary: str,
    ) -> None:
        """Add executive summary slide with readable narrative text."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_frame.text = "Executive Summary"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(44, 62, 80)
        
        # Summary content
        summary_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1), Inches(9), Inches(4.2)
        )
        summary_frame = summary_box.text_frame
        summary_frame.word_wrap = True
        
        # Parse paragraphs from the summary
        paragraphs = executive_summary.strip().split('\n\n')
        first_para = True
        
        for para_text in paragraphs:
            if not para_text.strip():
                continue
                
            # Clean up markdown formatting
            clean_text = para_text.strip()
            # Remove ** bold markers for PPTX
            import re
            clean_text = re.sub(r'\*\*(.+?)\*\*', r'\1', clean_text)
            
            if first_para:
                p = summary_frame.paragraphs[0]
                first_para = False
            else:
                p = summary_frame.add_paragraph()
            
            p.text = clean_text
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(52, 73, 94)
            p.space_after = Pt(12)

    def _add_metrics_slide(
        self, 
        prs: Presentation,
        code_result: dict[str, Any] | None,
        code_output: str | None
    ) -> None:
        """Add key metrics slide(s) with smart pagination."""
        if not code_result:
            # Fallback to code_output if no structured result
            if code_output:
                self._add_metrics_slide_from_output(prs, code_output)
            return
        
        # Split code_result into chunks that fit on slides
        # Each slide can comfortably fit ~8-10 main items (considering sub-bullets for dicts)
        items = list(code_result.items())
        slides_data = self._split_metrics_into_slides(items, max_items_per_slide=8)
        
        for slide_num, slide_items in enumerate(slides_data):
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
            
            # Slide title
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.4), Inches(9), Inches(0.6)
            )
            title_frame = title_box.text_frame
            if len(slides_data) > 1:
                title_frame.text = f"Key Findings (Part {slide_num + 1}/{len(slides_data)})"
            else:
                title_frame.text = "Key Findings"
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(32)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(44, 62, 80)
            
            # Content area
            content_box = slide.shapes.add_textbox(
                Inches(1), Inches(1.3), Inches(8), Inches(4.5)
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            
            # Add metrics for this slide
            self._add_metrics_to_frame(content_frame, slide_items)
    
    def _split_metrics_into_slides(self, items: list[tuple[str, Any]], max_items_per_slide: int = 8) -> list[list[tuple[str, Any]]]:
        """Split metrics into chunks that fit comfortably on slides."""
        slides = []
        current_slide = []
        current_weight = 0
        
        for key, value in items:
            # Calculate weight (dicts with sub-bullets count as 2-3 items)
            weight = 1
            if isinstance(value, dict) and len(value) > 0:
                # Dict will create sub-bullets, counts as multiple items
                weight = min(len(value), 5) // 2 + 1  # Each 2 dict items = 1 weight
            
            # If adding this would overflow, start new slide
            if current_slide and (current_weight + weight > max_items_per_slide):
                slides.append(current_slide)
                current_slide = []
                current_weight = 0
            
            current_slide.append((key, value))
            current_weight += weight
        
        # Add remaining items
        if current_slide:
            slides.append(current_slide)
        
        return slides if slides else [[]]
    
    def _add_metrics_to_frame(self, content_frame, items: list[tuple[str, Any]]) -> None:
        """Add metrics to a text frame with proper formatting."""
        for i, (key, value) in enumerate(items):
            if i > 0:
                para = content_frame.add_paragraph()
            else:
                para = content_frame.paragraphs[0]
            
            para.level = 0
            
            # Format value by unpacking data structures
            if isinstance(value, float):
                # Check if it looks like a percentage (between 0 and 1)
                if 0 <= value <= 1:
                    formatted_value = f"{value * 100:.1f}%"
                else:
                    formatted_value = f"{value:,.2f}"
            elif isinstance(value, int):
                formatted_value = f"{value:,}"
            elif isinstance(value, list):
                # Show actual list contents, limited to fit slide
                if len(value) > 0:
                    # Check if numeric or string list
                    if all(isinstance(x, (int, float)) for x in value[:4]):
                        items = [f"{x:,.2f}" if isinstance(x, float) else f"{x:,}" for x in value[:4]]
                    else:
                        items = [str(v)[:30] for v in value[:4]]
                    formatted_value = ", ".join(items)
                    if len(value) > 4:
                        formatted_value += f" ... (+{len(value) - 4} more)"
                else:
                    formatted_value = "(empty)"
            elif isinstance(value, dict):
                # Show actual dict breakdown with limited items
                if len(value) > 0:
                    # Create sub-bullets for dict items
                    metric_name = key.replace('_', ' ').title()
                    metric_name = metric_name.replace('Avg', 'Average').replace('Pct', 'Percent')
                    
                    para.text = f"{metric_name}:"
                    para.font.size = Pt(18)
                    para.font.bold = True
                    para.font.color.rgb = RGBColor(52, 73, 94)
                    para.space_after = Pt(6)
                    
                    # Add sub-bullets for dict entries (max 3 per dict on slide)
                    for j, (k, v) in enumerate(list(value.items())[:3]):
                        sub_para = content_frame.add_paragraph()
                        sub_para.level = 1  # Indent
                        
                        if isinstance(v, float):
                            v_str = f"{v:,.2f}"
                        elif isinstance(v, int):
                            v_str = f"{v:,}"
                        else:
                            v_str = str(v)[:30]
                        
                        sub_para.text = f"{k}: {v_str}"
                        sub_para.font.size = Pt(14)
                        sub_para.font.color.rgb = RGBColor(52, 73, 94)
                        sub_para.space_after = Pt(4)
                    
                    if len(value) > 3:
                        more_para = content_frame.add_paragraph()
                        more_para.level = 1
                        more_para.text = f"... (+{len(value) - 3} more)"
                        more_para.font.size = Pt(12)
                        more_para.font.italic = True
                        more_para.font.color.rgb = RGBColor(127, 140, 141)
                    continue  # Skip the standard formatting below
            else:
                formatted_value = str(value)[:100]  # Limit string length
            
            # Clean up the metric name
            metric_name = key.replace('_', ' ').title()
            # Handle common abbreviations
            metric_name = metric_name.replace('Avg', 'Average')
            metric_name = metric_name.replace('Pct', 'Percent')
            
            para.text = f"{metric_name}: {formatted_value}"
            para.font.size = Pt(16)  # Slightly smaller to fit more
            para.font.color.rgb = RGBColor(52, 73, 94)
            para.space_after = Pt(8)  # Reduced spacing
        
    
    def _add_metrics_slide_from_output(
        self, 
        prs: Presentation,
        code_output: str
    ) -> None:
        """Add metrics slide from code output when structured result unavailable."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Slide title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4), Inches(9), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_frame.text = "Key Findings"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(44, 62, 80)
        
        # Content area
        content_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.3), Inches(8), Inches(4.5)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        # Try to extract meaningful lines (avoid raw dict/json output)
        lines = [line.strip() for line in code_output.split('\n') if line.strip()]
        
        # Filter out lines that look like raw python dicts or JSON
        clean_lines = []
        for line in lines:
            # Skip lines that are clearly dict/json representations
            if line.startswith(('{', '}', '[', ']', "'", '"')):
                continue
            # Skip lines with too many brackets or quotes (likely raw data)
            if line.count('{') + line.count('[') + line.count("'") > 3:
                continue
            # Keep lines that look like actual insights (max 10 lines)
            if len(line) > 10 and len(line) < 200 and len(clean_lines) < 10:
                clean_lines.append(line)
        
        if clean_lines:
            for i, line in enumerate(clean_lines):
                if i > 0:
                    para = content_frame.add_paragraph()
                else:
                    para = content_frame.paragraphs[0]
                
                para.level = 0
                para.text = line
                para.font.size = Pt(16)
                para.font.color.rgb = RGBColor(52, 73, 94)
                para.space_after = Pt(10)
        else:
            # Fallback: show truncated output without raw data
            para = content_frame.paragraphs[0]
            para.text = "Analysis completed successfully. See visualizations for details."
            para.font.size = Pt(16)
            para.font.color.rgb = RGBColor(52, 73, 94)

    def _add_chart_slide(
        self, 
        prs: Presentation,
        img_bytes: bytes,
        title: str
    ) -> None:
        """Add a chart slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Slide title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(44, 62, 80)
        title_para.alignment = PP_ALIGN.CENTER
        
        # Add chart image (centered, large)
        try:
            img_stream = BytesIO(img_bytes)
            left = Inches(0.5)
            top = Inches(1.0)
            width = Inches(9)
            height = Inches(4.3)
            
            slide.shapes.add_picture(
                img_stream,
                left, top,
                width=width,
                height=height
            )
        except Exception as e:
            logger.error("Failed to add chart image: %s", e)
            # Add error text
            error_box = slide.shapes.add_textbox(
                Inches(2), Inches(2.5), Inches(6), Inches(1)
            )
            error_frame = error_box.text_frame
            error_para = error_frame.paragraphs[0]
            error_para.text = "[Chart could not be rendered]"
            error_para.font.size = Pt(16)
            error_para.font.italic = True
            error_para.font.color.rgb = RGBColor(192, 57, 43)
            error_para.alignment = PP_ALIGN.CENTER

    def _add_insights_slide(
        self, 
        prs: Presentation,
        code_output: str
    ) -> None:
        """Add insights/recommendations slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Slide title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4), Inches(9), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_frame.text = "Insights & Observations"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(44, 62, 80)
        
        # Content - extract meaningful insights from output
        content_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.3), Inches(8), Inches(3.8)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        # Try to extract clean, meaningful insights
        lines = [line.strip() for line in code_output.split('\n') if line.strip()]
        
        # Filter for presentation-worthy insights
        insights = []
        for line in lines:
            # Skip raw dict/json/code output
            if any(char in line for char in ['{', '}', '[', ']', "':", '":']):
                # But allow lines with minimal formatting characters
                if line.count('{') + line.count('[') + line.count("'") > 2:
                    continue
            
            # Skip very short or very long lines
            if len(line) < 15 or len(line) > 250:
                continue
            
            # Skip lines that look like code or variable assignments
            if '=' in line and not any(keyword in line.lower() for keyword in ['is', 'was', 'are', 'were']):
                continue
            
            insights.append(line)
        
        # If we found good insights, use them
        if insights:
            for i, insight in enumerate(insights[:5]):  # Max 5 insights
                if i > 0:
                    para = content_frame.add_paragraph()
                else:
                    para = content_frame.paragraphs[0]
                
                para.level = 0
                # Clean up the insight text
                cleaned = insight.strip('.,;:')
                # Capitalize first letter if not already
                if cleaned and cleaned[0].islower():
                    cleaned = cleaned[0].upper() + cleaned[1:]
                
                para.text = cleaned
                para.font.size = Pt(16)
                para.font.color.rgb = RGBColor(52, 73, 94)
                para.space_after = Pt(12)
        else:
            # Fallback: Show a professional summary message
            para = content_frame.paragraphs[0]
            para.text = "The analysis reveals important patterns in the data. Key trends and relationships are visualized in the charts above. Further investigation of specific segments may yield additional insights."
            para.font.size = Pt(16)
            para.font.color.rgb = RGBColor(52, 73, 94)
            para.alignment = PP_ALIGN.LEFT
